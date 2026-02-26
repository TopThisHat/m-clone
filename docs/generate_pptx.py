"""
Generates Playbook-Research-Architecture.pptx  — polished redesign.
Run: python3 docs/generate_pptx.py
"""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Design tokens ──────────────────────────────────────────────────────────────
BG     = RGBColor(0x0B, 0x16, 0x29)   # page background
S1     = RGBColor(0x0F, 0x20, 0x40)   # card surface
S2     = RGBColor(0x14, 0x29, 0x52)   # raised surface
BORD   = RGBColor(0x1E, 0x3A, 0x6E)   # subtle border
GOLD   = RGBColor(0xD4, 0xAA, 0x50)   # primary accent
GOLDL  = RGBColor(0xF0, 0xC9, 0x6B)   # lighter gold
WT     = RGBColor(0xF1, 0xF5, 0xF9)   # primary text
T2     = RGBColor(0x94, 0xA3, 0xB8)   # secondary text
T3     = RGBColor(0x47, 0x55, 0x69)   # muted text / footer
EME    = RGBColor(0x10, 0xB9, 0x81)   # emerald  — backend / data
SKY    = RGBColor(0x38, 0xBD, 0xF8)   # sky blue — frontend
AMB    = RGBColor(0xF5, 0x9E, 0x0B)   # amber    — agent
VIO    = RGBColor(0x8B, 0x5C, 0xF6)   # violet   — infra / auth
TEA    = RGBColor(0x06, 0xB6, 0xD4)   # teal     — AWS / secrets
ROS    = RGBColor(0xF4, 0x3F, 0x5E)   # rose     — warnings
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FT  = "Calibri Light"   # titles
FB  = "Calibri"         # body
FM  = "Consolas"        # code / mono

SW = Inches(13.33)
SH = Inches(7.5)
TOTAL = 16

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]
_n = [0]

# ── Primitives ─────────────────────────────────────────────────────────────────

def new_slide():
    _n[0] += 1
    return prs.slides.add_slide(blank)

def box(sl, l, t, w, h, fill, lc=None, lw=Pt(0.75)):
    sh = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if lc:  sh.line.color.rgb = lc; sh.line.width = lw
    else:   sh.line.fill.background()
    return sh

def t(sl, text, l, top, w, h, color=WT, size=11, bold=False,
      align=PP_ALIGN.LEFT, font=FB, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(Inches(l), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    r.font.italic = italic

def hline(sl, l, top, w, color=BORD, width=Pt(0.6)):
    """Thin horizontal rule."""
    sh = sl.shapes.add_connector(1, Inches(l), Inches(top),
                                 Inches(l + w), Inches(top))
    sh.line.color.rgb = color; sh.line.width = width

def chrome(sl, title, category="", accent=GOLD):
    """Standard slide chrome: background, header bar, separator, footer."""
    box(sl, 0, 0, 13.33, 7.5, BG)           # background
    box(sl, 0, 0, 13.33, 0.055, accent)      # top accent strip
    box(sl, 0, 0.055, 13.33, 0.88, S1)       # header surface
    hline(sl, 0, 0.935, 13.33, color=BORD)   # separator below header
    if category:
        t(sl, category.upper(), 0.45, 0.10, 9, 0.28,
          color=accent, size=8.5, bold=True, font=FB)
    t(sl, title, 0.45, 0.37, 11.5, 0.54,
      color=WHITE, size=27, bold=True, font=FT)
    # Footer
    hline(sl, 0, 7.17, 13.33, color=accent, width=Pt(0.4))
    box(sl, 0, 7.18, 13.33, 0.32, S1)
    t(sl, "PLAYBOOK RESEARCH", 0.45, 7.2, 3.5, 0.28,
      color=accent, size=8, bold=True)
    t(sl, "CONFIDENTIAL  ·  DO NOT DISTRIBUTE", 4.5, 7.2, 4.5, 0.28,
      color=T3, size=7.5, align=PP_ALIGN.CENTER)
    t(sl, str(_n[0]), 12.3, 7.2, 0.65, 0.28,
      color=T2, size=8, align=PP_ALIGN.RIGHT)

def card(sl, l, top, w, h, accent, title, items,
         ts=11.5, is_=10.5, fill=S1, icon=None):
    """Left-border accent card with title + bullet list."""
    box(sl, l, top, w, h, fill, lc=BORD, lw=Pt(0.5))
    box(sl, l, top, 0.045, h, accent)          # left accent strip
    x0 = l + 0.14
    if icon:
        t(sl, icon, x0, top + 0.08, 0.4, 0.32, color=accent, size=ts, bold=True)
        t(sl, title, x0 + 0.38, top + 0.08, w - 0.6, 0.32, color=accent, size=ts, bold=True)
    else:
        t(sl, title, x0, top + 0.08, w - 0.22, 0.32, color=accent, size=ts, bold=True)
    hline(sl, l + 0.14, top + 0.44, w - 0.22, color=BORD, width=Pt(0.4))
    row_h = max((h - 0.55) / max(len(items), 1), 0.3)
    for i, item in enumerate(items):
        t(sl, item, x0, top + 0.5 + i * row_h, w - 0.22, row_h,
          color=WT if i % 2 == 0 else T2, size=is_)

def step_badge(sl, l, top, num, color, sz=0.52):
    """Filled square badge with a number (simulates a circle badge)."""
    box(sl, l, top, sz, sz, color)
    t(sl, str(num), l, top + 0.02, sz, sz - 0.04,
      color=BG, size=18, bold=True, font=FT, align=PP_ALIGN.CENTER)

def tag(sl, l, top, w, h, text, fill, text_color=BG):
    box(sl, l, top, w, h, fill)
    t(sl, text, l, top, w, h, color=text_color, size=9,
      bold=True, align=PP_ALIGN.CENTER)

def method_pill(sl, l, top, method, path, desc, mc):
    tag(sl, l, top, 0.78, 0.38, method, mc)
    t(sl, path, l + 0.84, top + 0.01, 3.6, 0.22, color=GOLDL, size=8.5, bold=True, font=FM)
    t(sl, desc, l + 0.84, top + 0.21, 6.0, 0.18, color=T2, size=8)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, 13.33, 7.5, BG)
box(sl, 0, 0, 13.33, 0.055, GOLD)
box(sl, 0, 7.45, 13.33, 0.055, GOLD)

# Subtle diagonal accent block (top-right)
box(sl, 9.5, 0.055, 3.83, 7.4, S1)
hline(sl, 9.5, 0.055, 0, color=GOLD, width=Pt(1.5))   # left edge of accent block

# P Logo
box(sl, 1.1, 1.55, 1.55, 1.55, GOLD)
t(sl, "P", 1.1, 1.55, 1.55, 1.55, color=BG, size=72, bold=True, font=FT, align=PP_ALIGN.CENTER)

# Product name
t(sl, "Playbook Research", 1.1, 3.28, 8.2, 0.78,
  color=WHITE, size=46, bold=True, font=FT)
t(sl, "AI-Powered Deep Research Platform", 1.1, 4.08, 8.0, 0.45,
  color=GOLD, size=20, font=FT)
hline(sl, 1.1, 4.65, 7.5, color=BORD)
t(sl, "System Architecture Overview", 1.1, 4.75, 7.5, 0.38,
  color=T2, size=14, font=FB)

# Stats row
stats = [("7", "Agent Tools"), ("12", "DB Tables"), ("8", "API Routers"), ("14", "Frontend Routes")]
for i, (val, lbl) in enumerate(stats):
    x = 1.1 + i * 1.85
    box(sl, x, 5.4, 1.65, 0.95, S1, lc=BORD, lw=Pt(0.5))
    box(sl, x, 5.4, 1.65, 0.04, GOLD)
    t(sl, val, x, 5.45, 1.65, 0.5, color=GOLD, size=22, bold=True, font=FT, align=PP_ALIGN.CENTER)
    t(sl, lbl, x, 5.92, 1.65, 0.3, color=T2, size=9, align=PP_ALIGN.CENTER)

t(sl, "Confidential  ·  Version 1.0", 1.1, 6.55, 8.0, 0.3, color=T3, size=10)

# Accent column content
t(sl, "Stack", 9.75, 0.7, 3.2, 0.3, color=GOLD, size=9, bold=True)
hline(sl, 9.75, 1.05, 3.1, color=BORD)
stack = [
    (EME,  "FastAPI  ·  Python 3.12"),
    (SKY,  "SvelteKit  ·  TypeScript"),
    (AMB,  "PydanticAI  ·  GPT-4o"),
    (VIO,  "PostgreSQL  ·  asyncpg"),
    (TEA,  "Redis  ·  ElastiCache"),
    (GOLDL,"AWS Secrets Manager"),
]
for i, (c, s) in enumerate(stack):
    y = 1.15 + i * 0.82
    box(sl, 9.75, y, 0.06, 0.55, c)
    t(sl, s, 9.87, y + 0.04, 3.0, 0.5, color=WT, size=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PLATFORM OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "Platform Overview", "Introduction", GOLD)

pillars = [
    ("Research Engine", EME,
     ["Mandatory 4-phase agent loop enforces depth",
      "7 specialised tools: web, wiki, finance, SEC, docs",
      "Real-time SSE streaming — tokens as they're generated",
      "Self-evaluation & confidence scoring before writing",
      "Inline citations + numeric conflict detection"]),
    ("Team Collaboration", SKY,
     ["Create teams with role-based access (owner/admin/member)",
      "Share sessions: private → team → public",
      "Threaded comments with @mention notifications",
      "Pin sessions to team workspace",
      "Full team activity feed"]),
    ("Memory & Intelligence", AMB,
     ["Post-research entity extraction (GPT-4o-mini)",
      "Cross-session context enrichment",
      "User-defined domain research rules",
      "BM25 full-text search over uploaded PDFs",
      "Usage tracking & cost dashboard"]),
    ("Production Infrastructure", VIO,
     ["AWS Secrets Manager with live credential rotation",
      "PostgreSQL RDS + ElastiCache Redis (TLS)",
      "OIDC SSO — any OpenID Connect provider",
      "Docker Compose for zero-config local dev",
      "Async webhook jobs for external integrations"]),
]

for i, (title, color, items) in enumerate(pillars):
    col, row = i % 2, i // 2
    card(sl, 0.38 + col * 6.55, 1.02 + row * 2.95,
         6.32, 2.82, color, title, items, ts=12, is_=10.5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — USE CASES
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "Use Cases", "Who Uses Playbook Research & Why", GOLD)

use_cases = [
    ("Investment & Financial Research", GOLDL, "◆",
     ["Analyze company fundamentals, earnings, and ratios via Yahoo Finance",
      "Pull SEC 10-K / 10-Q / 8-K filings directly through EDGAR",
      "Cross-reference web sources with financial data for conflict detection",
      "Generate cite-backed reports ready to share with stakeholders"]),
    ("Competitive Intelligence", SKY, "⊙",
     ["Map competitor product lines, pricing, and market positioning",
      "Track industry news and analyst commentary in real time",
      "Upload competitor whitepapers or annual reports for BM25 search",
      "Share findings to a team workspace with pinned sessions"]),
    ("Due Diligence", VIO, "◈",
     ["Deep-dive research across web, wiki, filings, and uploaded docs",
      "Confidence scoring flags gaps before the report is written",
      "Inline citations with source URLs — every claim is traceable",
      "Audit trail preserved in session history for future reference"]),
    ("Industry & Market Analysis", EME, "◉",
     ["Broad multi-source synthesis: news, Wikipedia, financials, filings",
      "Minimum 4 tool calls per query ensures coverage breadth",
      "Numeric conflict detection flags diverging data across sources",
      "Export as a shareable read-only report at /share/{uuid}"]),
    ("Document Analysis", AMB, "▣",
     ["Upload PDFs — research papers, contracts, earnings transcripts",
      "BM25 full-text search surfaces the most relevant passages",
      "Agent combines uploaded doc context with live web research",
      "Useful for annotating internal reports with external validation"]),
    ("Team Knowledge Building", TEA, "◎",
     ["Share sessions to team workspaces; members view at /share/{id}",
      "Threaded comments with @mentions for async collaboration",
      "Pin the most valuable sessions to the top of the team dashboard",
      "Activity feed keeps everyone aligned on what's been researched"]),
]

for i, (title, color, icon, items) in enumerate(use_cases):
    col, row = i % 3, i // 3
    x = 0.38 + col * 4.3
    y = 1.02 + row * 3.04
    w, h = 4.18, 2.92
    box(sl, x, y, w, h, S1, lc=BORD, lw=Pt(0.5))
    box(sl, x, y, 0.045, h, color)
    # Icon badge
    box(sl, x + 0.14, y + 0.1, 0.44, 0.44, color)
    t(sl, icon, x + 0.14, y + 0.1, 0.44, 0.44,
      color=BG, size=14, bold=True, align=PP_ALIGN.CENTER)
    t(sl, title, x + 0.66, y + 0.16, w - 0.8, 0.32,
      color=color, size=11, bold=True)
    hline(sl, x + 0.14, y + 0.6, w - 0.22, color=BORD, width=Pt(0.4))
    for j, item in enumerate(items):
        t(sl, item, x + 0.18, y + 0.68 + j * 0.5, w - 0.3, 0.46,
          color=WT if j % 2 == 0 else T2, size=9.5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW TO USE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "How to Use Playbook Research", "User Guide · End-to-End Workflow", SKY)

# ── Top: linear user journey (8 steps across the width)
steps = [
    ("Log In",       SKY,  "Sign in via OIDC SSO or use\ndev-bypass locally"),
    ("Ask",          GOLDL,"Type your research question\nin the chat input"),
    ("Configure",    AMB,  "Choose model (GPT-4o / Claude)\nand depth (fast · balanced · deep)"),
    ("Upload",       VIO,  "Optionally attach PDFs —\nearnings reports, whitepapers"),
    ("Watch",        EME,  "See tool calls + reasoning\nstream live as SSE events"),
    ("Review",       GOLDL,"Read the final report with\ninline citations and sources"),
    ("Share",        SKY,  "Share to your team or make\npublic via /share/{uuid}"),
    ("Collaborate",  TEA,  "Comment, @mention teammates,\npin to team workspace"),
]

step_w = 1.55
for i, (label, color, desc) in enumerate(steps):
    x = 0.38 + i * 1.57
    # Connector
    if i < len(steps) - 1:
        box(sl, x + step_w, 1.47, 0.02, 0.18, T3)
    # Number badge
    box(sl, x, 1.3, step_w, 0.52, color)
    t(sl, f"{i + 1}  {label}", x + 0.06, 1.32, step_w - 0.1, 0.48,
      color=BG, size=10, bold=True, font=FT)
    # Description
    box(sl, x, 1.85, step_w, 0.92, S1, lc=color, lw=Pt(0.8))
    t(sl, desc, x + 0.08, 1.92, step_w - 0.12, 0.82,
      color=WT, size=9, wrap=True)

hline(sl, 0.38, 2.9, 12.57, color=BORD)

# ── Bottom: three detailed user scenarios
scenarios = [
    ("Analyst: Solo Research Session", GOLDL, [
        ('Enter query', "e.g. 'Summarise Tesla Q4 2024 earnings and analyst reactions'"),
        ('Set depth',   "Choose 'deep' for maximum source coverage (6+ tool calls)"),
        ('Upload docs', "Drop in the earnings PDF — agent will BM25-search it alongside web"),
        ('Stream',      "Watch web_search, get_financials, wiki_lookup fire in the swimlane"),
        ('Export',      "Copy the markdown report or share the permanent /share/{uuid} link"),
    ]),
    ("Team Lead: Collaborative Research", SKY, [
        ('Create team', "Settings → Teams → Create → invite colleagues by email"),
        ('Research',    "Run the session as usual; the report is saved automatically"),
        ('Share',       "Session card menu → Share to Team → select your team"),
        ('Notify',      "Team members see it in the team dashboard and activity feed"),
        ('Discuss',     "@mention teammates in comments; they get an inbox notification"),
    ]),
    ("Developer: Async / Webhook Integration", TEA, [
        ('POST job',    "POST /api/research/async  {query, webhook_url, model, depth}"),
        ('Poll',        "GET /api/research/jobs/{id}  →  {status: queued|running|done}"),
        ('Receive',     "On completion your webhook_url receives the full result_markdown"),
        ('Embed',       "Use GET /api/share/{id} to fetch the structured JSON for your app"),
        ('Automate',    "Chain queries, set schedules, or trigger on upstream data events"),
    ]),
]

for i, (title, color, steps_) in enumerate(scenarios):
    x = 0.38 + i * 4.32
    w, h = 4.18, 3.95
    box(sl, x, 3.05, w, h, S1, lc=BORD, lw=Pt(0.5))
    box(sl, x, 3.05, 0.045, h, color)
    t(sl, title, x + 0.14, 3.12, w - 0.2, 0.3, color=color, size=11, bold=True)
    hline(sl, x + 0.14, 3.46, w - 0.2, color=BORD, width=Pt(0.4))
    for j, (action, detail) in enumerate(steps_):
        y = 3.54 + j * 0.7
        box(sl, x + 0.14, y, 0.22, 0.22, color)
        t(sl, str(j + 1), x + 0.14, y, 0.22, 0.22,
          color=BG, size=8, bold=True, align=PP_ALIGN.CENTER)
        t(sl, action, x + 0.42, y, 0.85, 0.22, color=color, size=9, bold=True)
        t(sl, detail,  x + 0.42, y + 0.24, w - 0.6, 0.38, color=T2, size=8.8, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "System Architecture", "Overview", GOLD)

diagram = os.path.join(os.path.dirname(os.path.abspath(__file__)), "architecture-diagram.png")
if os.path.exists(diagram):
    sl.shapes.add_picture(diagram, Inches(0.25), Inches(0.98),
                          Inches(12.83), Inches(6.12))
else:
    t(sl, "Run generate_diagram.py first.", 2, 3.5, 9, 0.5,
      color=ROS, size=16, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RESEARCH AGENT
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "Research Agent — Mandatory 4-Phase Loop", "Backend · PydanticAI", AMB)

phases = [
    ("0", "PLAN",     AMB,  "create_research_plan\ncalled first — always.\nScopes questions\nand source strategy."),
    ("1", "EXECUTE",  EME,  "Minimum 4 tool calls\nwith distinct queries:\nweb · wiki · finance\ndocs · SEC EDGAR"),
    ("2", "EVALUATE", SKY,  "evaluate_research_\ncompleteness scores\nbreadth, depth,\nrecency 0–100"),
    ("3", "DEEPER",   GOLDL,"If confidence < 85%:\ntargeted follow-up\nsearches fill\nidentified gaps"),
    ("4", "REPORT",   VIO,  "Final markdown with\ninline citations.\nConflict warnings\nif >25% variance"),
]

gap = 2.48
for i, (num, label, color, desc) in enumerate(phases):
    x = 0.38 + i * gap
    # Connector arrow (except last)
    if i < 4:
        box(sl, x + 2.04, 1.72, 0.44, 0.06, T3)   # arrow body
        box(sl, x + 2.42, 1.65, 0.06, 0.2, T3)     # arrowhead right side
    # Phase box
    box(sl, x, 1.02, 2.0, 0.72, color)
    t(sl, f"Phase {num}  ·  {label}", x + 0.08, 1.08, 1.85, 0.28,
      color=BG, size=11.5, bold=True, font=FT)
    # Detail card
    box(sl, x, 1.78, 2.0, 2.55, S1, lc=color, lw=Pt(1.0))
    box(sl, x, 1.78, 2.0, 0.04, color)
    t(sl, desc, x + 0.1, 1.88, 1.82, 2.38,
      color=WT, size=10.5, wrap=True)

# SSE Events strip
hline(sl, 0.38, 4.5, 12.57, color=BORD)
t(sl, "Real-time SSE Events", 0.38, 4.6, 3.0, 0.28, color=GOLD, size=10, bold=True)
events = ["start", "tool_call_start", "tool_executing", "tool_result",
          "text_delta", "chart_data", "conflict_warning", "final_report", "done"]
for i, ev in enumerate(events):
    x = 0.38 + i * 1.38
    box(sl, x, 5.0, 1.28, 0.32, S2, lc=AMB, lw=Pt(0.6))
    t(sl, ev, x + 0.06, 5.01, 1.18, 0.3, color=AMB, size=8.5, bold=True)

t(sl, "Agent streams each event type progressively — browser receives tokens and tool steps in real time without polling.",
  0.38, 5.45, 12.57, 0.3, color=T2, size=9.5, italic=True)

# Model options
t(sl, "Supported models:", 0.38, 5.9, 2.0, 0.28, color=T2, size=9.5)
for i, m in enumerate(["OpenAI GPT-4o (default)", "Anthropic Claude Sonnet 4.6"]):
    box(sl, 2.4 + i * 4.5, 5.87, 4.2, 0.3, S2, lc=BORD, lw=Pt(0.5))
    t(sl, m, 2.5 + i * 4.5, 5.88, 4.0, 0.28, color=WT, size=9.5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AGENT TOOLS
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "Agent Tools", "Backend · Research Capabilities", AMB)

tools = [
    ("create_research_plan",        AMB,  "◎", "Structures investigation scope, key questions, and source strategy before any research begins. Called first — mandatory."),
    ("evaluate_research_completeness", ROS, "◈", "Self-assessment scores breadth, depth, recency, and source diversity 0–100. Triggers dig-deeper phase if below 85%."),
    ("web_search",                  EME,  "⊙", "Tavily API integration. Returns up-to-date web results with snippets. All source URLs tracked for citation validation."),
    ("wiki_lookup",                 SKY,  "◉", "Wikipedia API for encyclopedic background on companies, industries, concepts, and historical context."),
    ("get_financials",              GOLDL,"◆", "Yahoo Finance (yfinance). Price history, fundamentals, earnings, ratios. Chart payloads streamed live to the browser."),
    ("search_uploaded_documents",   VIO,  "▣", "BM25 Okapi full-text search over user-uploaded PDFs. Text extracted by pdfplumber, cached in Redis with 24h TTL."),
    ("sec_edgar_search",            TEA,  "◧", "SEC EDGAR filing lookup. 10-K, 10-Q, 8-K for public company regulatory filings and official disclosures."),
]

cols = 2
for i, (name, color, icon, desc) in enumerate(tools):
    col = i % cols
    row = i // cols
    x = 0.38 + col * 6.55
    y = 1.02 + row * 1.48
    h = 1.35
    box(sl, x, y, 6.32, h, S1, lc=BORD, lw=Pt(0.5))
    box(sl, x, y, 0.045, h, color)
    t(sl, icon, x + 0.12, y + 0.07, 0.42, 0.38, color=color, size=18)
    t(sl, name, x + 0.58, y + 0.08, 5.55, 0.3, color=color, size=11, bold=True, font=FM)
    hline(sl, x + 0.14, y + 0.42, 6.0, color=BORD, width=Pt(0.4))
    t(sl, desc, x + 0.14, y + 0.5, 6.0, 0.75, color=T2, size=10, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DATA MODEL
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "Database Schema", "Data · PostgreSQL · 12 Tables", VIO)

groups = [
    ("Core", VIO, [
        ("sessions",  ["id UUID (PK)", "title · query", "report_markdown (TEXT)", "message_history / trace_steps (JSONB)", "owner_sid → users", "visibility: private | team | public"]),
        ("users",     ["sid (PK)", "display_name · email", "avatar_url · theme", "created_at · last_login"]),
        ("agent_memory",["session_id → sessions", "entity · entity_type", "facts (JSONB array)"]),
    ]),
    ("Collaboration", SKY, [
        ("teams",      ["id UUID (PK)", "slug (UNIQUE)", "display_name · description", "created_by → users"]),
        ("team_members",["team_id + sid (composite PK)", "role: owner | admin | member | viewer", "joined_at"]),
        ("session_teams",["session_id + team_id (PK)", "shared_at — many-to-many"]),
    ]),
    ("Engagement", EME, [
        ("comments",   ["session_id → sessions", "author_sid → users", "body · mentions (JSONB)", "parent_id (threading)"]),
        ("pinned_sessions",["sid + session_id + team_id (PK)", "pinned_at"]),
        ("notifications",["recipient_sid → users", "type · payload (JSONB)", "read (BOOLEAN)"]),
    ]),
    ("Operations", AMB, [
        ("team_activity",["team_id → teams", "actor_sid → users", "action · payload (JSONB)", "created_at"]),
        ("research_jobs",["id UUID (PK)", "query · webhook_url", "status: queued|running|done|failed", "result_markdown · error"]),
    ]),
]

col_x = [0.38, 3.72, 7.05, 10.38]
for gi, (group_name, gcolor, tables) in enumerate(groups):
    x = col_x[gi]
    w = 3.18
    box(sl, x, 1.02, w, 0.36, gcolor)
    t(sl, group_name.upper(), x, 1.02, w, 0.36,
      color=BG, size=10.5, bold=True, align=PP_ALIGN.CENTER)
    y = 1.42
    for tname, fields in tables:
        th = 0.34 + len(fields) * 0.3
        box(sl, x, y, w, th, S1, lc=BORD, lw=Pt(0.5))
        box(sl, x, y, w, 0.3, S2)
        box(sl, x, y, 0.04, th, gcolor)
        t(sl, tname, x + 0.1, y + 0.03, w - 0.15, 0.26,
          color=gcolor, size=10, bold=True, font=FM)
        for fi, f in enumerate(fields):
            t(sl, f, x + 0.1, y + 0.33 + fi * 0.3, w - 0.15, 0.28, color=T2, size=8.5)
        y += th + 0.1


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — COLLABORATION MODEL
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "Session Sharing & Team RBAC", "Collaboration Model", SKY)

# ── Left: Visibility funnel
t(sl, "Session Visibility", 0.38, 1.08, 6.0, 0.32, color=SKY, size=13, bold=True)
hline(sl, 0.38, 1.44, 6.0, color=BORD)

vis = [
    ("private",  T3,   S1,  "Owner only. Default state for all new sessions."),
    ("team",     GOLDL, S2, "Shared to one or more teams. Accessible to any\nmember via /share/{uuid} with JWT auth."),
    ("public",   EME,  S2,  "is_public = TRUE. Anyone can read via /share/{uuid}\nwithout authentication. Permanent URL."),
]
for i, (state, sc, bg_, desc) in enumerate(vis):
    y = 1.55 + i * 1.6
    w = 5.7 - i * 0.5
    x = 0.38 + i * 0.25
    box(sl, x, y, w, 1.42, bg_, lc=sc, lw=Pt(1.2))
    box(sl, x, y, w, 0.34, sc)
    t(sl, state.upper(), x + 0.12, y + 0.04, w - 0.2, 0.28,
      color=BG if sc != T3 else WT, size=11, bold=True)
    t(sl, desc, x + 0.12, y + 0.42, w - 0.2, 0.9, color=WT, size=10, wrap=True)
    if i < 2:
        t(sl, "↓", x + w / 2 - 0.15, y + 1.44, 0.3, 0.2, color=T3, size=13, align=PP_ALIGN.CENTER)

# Share endpoints
t(sl, "POST /api/sessions/{id}/share  →  is_public=TRUE",
  0.38, 6.45, 6.0, 0.25, color=T2, size=9, font=FM)
t(sl, "POST /api/sessions/{id}/teams  →  visibility='team'",
  0.38, 6.72, 6.0, 0.25, color=T2, size=9, font=FM)

# ── Right: RBAC
hline(sl, 6.8, 1.02, 0, color=BORD)   # visual divider
t(sl, "Team Role Hierarchy", 7.0, 1.08, 6.0, 0.32, color=SKY, size=13, bold=True)
hline(sl, 7.0, 1.44, 6.0, color=BORD)

roles = [
    ("OWNER",  GOLD,  "Full control: edit team settings, manage all members,\ndelete or unshare any session."),
    ("ADMIN",  AMB,   "Manage members, share/unshare sessions, moderate\ncomments. Cannot delete team."),
    ("MEMBER", EME,   "View team sessions, share own sessions to the team,\npost comments and @mentions."),
    ("VIEWER", SKY,   "Read-only access. Can view shared sessions and\ncomments. No write access."),
]
for i, (role, color, desc) in enumerate(roles):
    y = 1.55 + i * 1.36
    w = 5.0 - i * 0.35
    box(sl, 7.0, y, w, 1.2, S1, lc=color, lw=Pt(1.2))
    box(sl, 7.0, y, 1.1, 1.2, color)
    t(sl, role, 7.0, y, 1.1, 1.2, color=BG, size=10.5, bold=True,
      align=PP_ALIGN.CENTER, font=FT)
    t(sl, desc, 8.18, y + 0.22, w - 1.28, 0.8, color=WT, size=10, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FRONTEND
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "Frontend Architecture", "SvelteKit · TypeScript", SKY)

# Route tree
t(sl, "Route Structure", 0.38, 1.08, 5.8, 0.3, color=SKY, size=12, bold=True)
hline(sl, 0.38, 1.42, 5.8, color=BORD)

routes = [
    (0, "+layout.server.ts",  GOLD,  "JWT guard · user hydration · redirect to /login"),
    (1, "+layout.svelte",     SKY,   "App shell: sidebar, nav, notification badge"),
    (1, "/(app)/+page.svelte",SKY,   "Research workspace — chat + session list"),
    (2, "/teams/[slug]",      GOLDL, "Team dashboard, members, pinned sessions"),
    (2, "/notifications",     VIO,   "Notification inbox (polled every 30 s)"),
    (0, "/share/[id]",        EME,   "Read-only report — team members + public"),
    (0, "/login",             T2,    "OIDC redirect trigger"),
    (0, "/dashboard",         AMB,   "Usage stats: tokens, cost, top queries"),
]
for i, (indent, path, color, desc) in enumerate(routes):
    y = 1.52 + i * 0.57
    if indent > 0:
        box(sl, 0.38 + (indent - 1) * 0.3, y + 0.13, 0.2, 0.04, BORD)
        box(sl, 0.38 + indent * 0.3, y + 0.04, 0.04, 0.26, BORD)
    box(sl, 0.38 + indent * 0.3, y, 3.5 - indent * 0.18, 0.46, S1, lc=BORD, lw=Pt(0.5))
    box(sl, 0.38 + indent * 0.3, y, 0.04, 0.46, color)
    t(sl, path, 0.52 + indent * 0.3, y + 0.07, 2.85, 0.28, color=color, size=9.5, bold=True, font=FM)
    t(sl, desc, 4.1, y + 0.07, 2.1, 0.28, color=T2, size=9)

# Key components
t(sl, "Key Components & Patterns", 6.5, 1.08, 6.5, 0.3, color=SKY, size=12, bold=True)
hline(sl, 6.5, 1.42, 6.5, color=BORD)

comps = [
    (AMB,  "ChatThread.svelte",
     "SSE EventSource listener. Streams tool_call_start, text_delta, final_report events and updates reactive stores in real time."),
    (GOLDL,"ResearchSwimlane.svelte",
     "Visual timeline of agent tool calls. Shows which tool ran, arguments passed, and a 400-char result preview per step."),
    (EME,  "api/sessions.ts",
     "Typed REST client. All fetch() calls use credentials: 'include' so the JWT cookie is forwarded — including on SSR runs via SvelteKit's enhanced fetch."),
    (VIO,  "Auth guard",
     "+layout.server.ts checks jwt cookie on every request. PUBLIC_PATHS=['/login','/share','/auth'] are exempt. Authenticated user injected into all page data."),
    (TEA,  "Notification polling",
     "setInterval 30 s → GET /api/notifications. Unread count drives the badge. PATCH /api/notifications/{id}/read on open."),
]
for i, (color, name, desc) in enumerate(comps):
    y = 1.52 + i * 1.1
    box(sl, 6.5, y, 6.5, 1.0, S1, lc=BORD, lw=Pt(0.5))
    box(sl, 6.5, y, 0.045, 1.0, color)
    t(sl, name, 6.62, y + 0.06, 6.2, 0.28, color=color, size=10.5, bold=True, font=FM)
    t(sl, desc, 6.62, y + 0.38, 6.2, 0.56, color=T2, size=9.5, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — AWS INFRASTRUCTURE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "AWS Infrastructure & Credential Rotation", "Infrastructure · Cloud", TEA)

# Left: Secret shapes
t(sl, "Secrets Manager — Secret Shape", 0.38, 1.08, 5.8, 0.3, color=TEA, size=12, bold=True)
hline(sl, 0.38, 1.42, 5.8, color=BORD)

secrets = [
    ("DB Secret  (AWS_SECRET_NAME)", VIO, [
        '  "host":     "mydb.cluster.us-east-1.rds.amazonaws.com"',
        '  "port":     5432',
        '  "username": "appuser"',
        '  "password": "s3cr3t!@#"   // percent-encoded in DSN',
        '  "dbname":   "research"',
    ]),
    ("Redis Secret  (AWS_ELASTICACHE_SECRET_NAME)", ROS, [
        '  "url":        "rediss://cluster.cache.amazonaws.com:6379"',
        '  "auth_token": "mytoken"    // passed as password param',
    ]),
]
y = 1.52
for title, color, lines in secrets:
    h = 0.38 + len(lines) * 0.36
    box(sl, 0.38, y, 5.85, h, S1, lc=color, lw=Pt(1.0))
    box(sl, 0.38, y, 0.045, h, color)
    t(sl, title, 0.5, y + 0.05, 5.6, 0.28, color=color, size=10, bold=True)
    box(sl, 0.38, y + 0.36, 5.85, h - 0.36, S2)
    t(sl, "{", 0.52, y + 0.36, 5.5, 0.28, color=T2, size=9.5, font=FM)
    for i, line_ in enumerate(lines):
        t(sl, line_, 0.52, y + 0.62 + i * 0.36, 5.55, 0.34, color=GOLDL, size=9, font=FM)
    t(sl, "}", 0.52, y + 0.62 + len(lines) * 0.36, 5.5, 0.28, color=T2, size=9.5, font=FM)
    y += h + 0.25

# Right: Rotation flow
t(sl, "Automatic Rotation Recovery", 6.6, 1.08, 6.5, 0.3, color=TEA, size=12, bold=True)
hline(sl, 6.6, 1.42, 6.5, color=BORD)

steps = [
    (TEA,  "AWS rotates the secret on schedule or on-demand"),
    (T3,   "Old password still valid during rotation window"),
    (ROS,  "pool.acquire() → new DB connection → auth rejected"),
    (AMB,  "asyncpg raises  InvalidPasswordError  (SQLSTATE 28P01)"),
    (GOLDL,"_acquire() catches → calls  _reset_pool()"),
    (EME,  "Pool closed · secret cache evicted (invalidate_db_secret)"),
    (TEA,  "get_db_secret() → Secrets Manager → fresh credentials"),
    (EME,  "New pool created with new DSN + ssl='require'"),
    (WT,   "Request retried — fully transparent to the caller"),
]
for i, (color, step) in enumerate(steps):
    y = 1.52 + i * 0.58
    box(sl, 6.6, y, 0.44, 0.44, color)
    t(sl, str(i + 1), 6.6, y, 0.44, 0.44, color=BG, size=12, bold=True,
      font=FT, align=PP_ALIGN.CENTER)
    box(sl, 7.1, y, 5.9, 0.44, S1, lc=BORD, lw=Pt(0.5))
    t(sl, step, 7.2, y + 0.07, 5.75, 0.3, color=WT, size=10)

t(sl, "Redis follows the same pattern — AuthenticationError → _reset_client() → invalidate_redis_secret() → re-fetch.",
  6.6, 6.82, 6.5, 0.28, color=T2, size=9, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DATA FLOWS
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "Key Data Flows", "Runtime Behaviour", EME)

flows = [
    ("Research Execution", EME, [
        "User submits query → POST /api/research",
        "Redis: load cached PDF context for session",
        "DB: fetch agent_memory for context enrichment",
        "Agent runs 4-phase loop (plan → execute → evaluate → report)",
        "SSE events stream tokens + tool calls live to browser",
        "final_report event → client persists session to DB",
        "Background: entity extraction → agent_memory (non-blocking)",
    ]),
    ("Comment & Notification", AMB, [
        "User posts comment with @username in body",
        "Backend regex-parses all @mention tokens",
        "Notification row inserted per mentioned user",
        "Frontend polls GET /api/notifications every 30 s",
        "Unread badge increments on the notification icon",
        "User opens inbox → PATCH /{id}/read or POST /read-all",
    ]),
    ("PDF Document Search", SKY, [
        "User uploads PDF → POST /api/documents/upload",
        "pdfplumber extracts full text → stored in Redis (24h TTL)",
        "session_key returned; included in subsequent research call",
        "search_uploaded_documents tool runs BM25 over extracted text",
        "Top-scoring chunks injected into agent context window",
    ]),
    ("Async Webhook Job", VIO, [
        "POST /api/research/async  {query, webhook_url}",
        "Job row created (status: queued) → job_id returned",
        "FastAPI BackgroundTask runs full agent loop async",
        "On completion: POST result_markdown to webhook_url",
        "Caller polls GET /api/research/jobs/{id} for status",
        "Status lifecycle: queued → running → done | failed",
    ]),
]

for i, (title, color, steps) in enumerate(flows):
    col, row = i % 2, i // 2
    x = 0.38 + col * 6.55
    y = 1.02 + row * 3.1
    h = 3.0
    box(sl, x, y, 6.32, h, S1, lc=BORD, lw=Pt(0.5))
    box(sl, x, y, 0.045, h, color)
    t(sl, title, x + 0.14, y + 0.08, 6.0, 0.3, color=color, size=12, bold=True)
    hline(sl, x + 0.14, y + 0.42, 6.0, color=BORD, width=Pt(0.4))
    for si, step in enumerate(steps):
        sy = y + 0.5 + si * 0.38
        t(sl, f"{si + 1}.", x + 0.14, sy, 0.3, 0.35, color=color, size=9.5, bold=True)
        t(sl, step, x + 0.44, sy, 5.7, 0.35, color=WT if si % 2 == 0 else T2, size=9.5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — API REFERENCE
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "API Reference", "Backend · FastAPI", GOLD)

MC = {"GET": EME, "POST": SKY, "PATCH": AMB, "DELETE": ROS}
endpoints = [
    # col 1
    ("POST",   "/api/research",                "Stream research (SSE)  ·  model, query, depth, pdf_session_key"),
    ("POST",   "/api/research/async",          "Async job + webhook  ·  returns job_id"),
    ("GET",    "/api/research/jobs/{id}",      "Poll async job status and result"),
    ("GET",    "/api/sessions",                "List sessions owned by current user"),
    ("POST",   "/api/sessions",                "Create new session"),
    ("GET",    "/api/sessions/{id}",           "Fetch session (owner / team check)"),
    ("PATCH",  "/api/sessions/{id}",           "Update title, report, visibility"),
    ("DELETE", "/api/sessions/{id}",           "Delete session"),
    ("POST",   "/api/sessions/{id}/share",     "Make public  →  returns share_url"),
    ("DELETE", "/api/sessions/{id}/share",     "Revoke public access"),
    # col 2
    ("POST",   "/api/sessions/{id}/teams",     "Share to team  →  visibility=team"),
    ("DELETE", "/api/sessions/{id}/teams/{t}", "Unshare from team"),
    ("GET",    "/api/share/{id}",              "Public/team read-only (auth optional)"),
    ("POST",   "/api/documents/upload",        "Upload PDF → Redis cache → session_key"),
    ("GET",    "/api/teams",                   "List user's teams"),
    ("POST",   "/api/teams",                   "Create team (caller becomes owner)"),
    ("GET",    "/api/teams/{slug}",            "Get team detail + member list"),
    ("PATCH",  "/api/teams/{slug}/members/{s}","Update member role"),
    ("GET",    "/api/notifications",           "Unread notification list"),
    ("GET",    "/health",                      "Health check  →  {status: ok}"),
]

split = 10
for i, (method, path, desc) in enumerate(endpoints):
    col = i // split
    row = i % split
    x = 0.38 + col * 6.52
    y = 1.02 + row * 0.6
    method_pill(sl, x, y, method, path, desc, MC.get(method, T3))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ARCHITECTURAL DECISIONS
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "Key Architectural Decisions", "Design Rationale", GOLD)

decisions = [
    ("SSE over WebSocket", GOLDL,
     ["One-way streaming maps perfectly to research output",
      "Works through all proxies and CDNs without upgrade headers",
      "Native EventSource API — zero extra library on client",
      "Compatible with serverless (no persistent connection needed)"]),
    ("asyncpg Direct — No ORM", EME,
     ["Maximum async performance — zero sync blocking overhead",
      "Advisory locks for safe concurrent schema migrations",
      "Native JSONB for message_history and trace_steps",
      "Simple parameterised queries without abstraction tax"]),
    ("Mandatory 4-Phase Loop", AMB,
     ["Prevents shallow single-tool responses",
      "Forces self-evaluation before writing begins",
      "Minimum 4 tool calls ensures research breadth",
      "Confidence scoring and conflict detection built in"]),
    ("Graceful Degradation", SKY,
     ["Redis unavailable → in-memory fallback, app keeps running",
      "DB not configured → endpoints return empty lists (no 500s)",
      "OIDC not set → dev bypass for friction-free local dev",
      "Each service fails independently — not catastrophically"]),
    ("Runtime Secret Rotation", TEA,
     ["No restart required when AWS rotates credentials",
      "Auth errors detected on next pool.acquire() call",
      "Single retry after re-fetching from Secrets Manager",
      "asyncio.Lock prevents thundering-herd on pool reset"]),
    ("3-Tier Visibility Model", VIO,
     ["private → team → public: clear escalation path",
      "Permanent share URLs: /share/{uuid} never changes",
      "Team membership verified server-side on every request",
      "4-level RBAC: owner > admin > member > viewer"]),
]

for i, (title, color, points) in enumerate(decisions):
    col, row = i % 3, i // 3
    card(sl, 0.38 + col * 4.3, 1.02 + row * 3.04,
         4.18, 2.92, color, title, points, ts=11.5, is_=10.5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — DEPLOYMENT
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, "Deployment Guide", "Infrastructure · Getting Started", EME)

# Local dev column
box(sl, 0.38, 1.02, 6.0, 0.38, EME)
t(sl, "Local Development", 0.52, 1.04, 5.7, 0.32, color=BG, size=12, bold=True)

local_steps = [
    ("1", EME,  "docker compose up -d",               "Starts Postgres 16 + Redis 7 with health checks"),
    ("2", GOLDL,"cp backend/.env.example backend/.env","Fill in OPENAI_API_KEY, TAVILY_API_KEY, JWT_SECRET"),
    ("3", SKY,  "DATABASE_URL=postgresql://...",       "localhost:5432/research (Docker default)"),
    ("4", SKY,  "REDIS_URL=redis://localhost:6379",    ""),
    ("5", VIO,  "DEV_AUTH_BYPASS=true",               "Skip OIDC — any username logs in during dev"),
    ("6", EME,  "uvicorn app.main:app --reload",       "Backend on :8000  (from backend/ directory)"),
    ("7", AMB,  "npm run dev",                         "Frontend on :5173  (from frontend/ directory)"),
]
for i, (num, color, cmd, desc) in enumerate(local_steps):
    y = 1.48 + i * 0.73
    box(sl, 0.38, y, 0.38, 0.55, color)
    t(sl, num, 0.38, y + 0.03, 0.38, 0.5, color=BG, size=12, bold=True,
      font=FT, align=PP_ALIGN.CENTER)
    box(sl, 0.8, y, 5.58, 0.55, S1, lc=BORD, lw=Pt(0.5))
    t(sl, cmd, 0.9, y + 0.03, 5.4, 0.26, color=GOLDL, size=9.5, bold=True, font=FM)
    if desc:
        t(sl, desc, 0.9, y + 0.3, 5.4, 0.22, color=T2, size=9)

# AWS column
box(sl, 6.95, 1.02, 6.0, 0.38, TEA)
t(sl, "AWS Production", 7.09, 1.04, 5.7, 0.32, color=BG, size=12, bold=True)

aws_vars = [
    ("AWS_SECRET_NAME",              "prod/myapp/db",    VIO,  "Postgres credentials (replaces DATABASE_URL)"),
    ("AWS_ELASTICACHE_SECRET_NAME",  "prod/myapp/redis", ROS,  "Redis token (replaces REDIS_URL)"),
    ("AWS_REGION",                   "us-east-1",        TEA,  "Shared by both secrets"),
    ("OPENAI_API_KEY",               "sk-...",           EME,  "GPT-4o research + entity extraction"),
    ("TAVILY_API_KEY",               "tvly-...",         AMB,  "Web search tool"),
    ("JWT_SECRET",                   "<strong key>",     GOLDL,"Change from default in all non-dev envs"),
    ("OIDC_ISSUER",                  "https://...",      SKY,  "Any OpenID Connect provider"),
    ("ALLOWED_ORIGINS",              '["https://..."]',  T2,   "CORS whitelist (JSON array)"),
]
for i, (key, val, color, desc) in enumerate(aws_vars):
    y = 1.48 + i * 0.72
    box(sl, 6.95, y, 6.0, 0.62, S1, lc=BORD, lw=Pt(0.5))
    box(sl, 6.95, y, 0.04, 0.62, color)
    t(sl, key, 7.06, y + 0.03, 3.4, 0.26, color=color, size=9, bold=True, font=FM)
    t(sl, val, 7.06, y + 0.32, 3.4, 0.24, color=GOLDL, size=8.5, font=FM)
    t(sl, desc, 10.55, y + 0.12, 2.25, 0.4, color=T2, size=8.5, wrap=True)

t(sl, "SSL is auto-enabled for Postgres; TLS (rediss://) is auto-enabled for Redis when using AWS secrets.",
  6.95, 7.05, 6.0, 0.24, color=T2, size=8.5, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — QUESTIONS
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, 13.33, 7.5, BG)
box(sl, 0, 0, 13.33, 0.055, GOLD)
box(sl, 0, 7.45, 13.33, 0.055, GOLD)

# Right accent
box(sl, 9.8, 0.055, 3.53, 7.4, S1)

# Large P
box(sl, 1.2, 1.6, 1.55, 1.55, GOLD)
t(sl, "P", 1.2, 1.6, 1.55, 1.55, color=BG, size=72, bold=True, font=FT, align=PP_ALIGN.CENTER)

t(sl, "Playbook Research", 1.2, 3.38, 8.0, 0.72, color=WHITE, size=42, bold=True, font=FT)
t(sl, "Questions?", 1.2, 4.12, 8.0, 0.52, color=GOLD, size=24, font=FT)
hline(sl, 1.2, 4.75, 7.0, color=BORD)

t(sl, "Thank you for your time.", 1.2, 4.88, 7.0, 0.38, color=T2, size=14)

summary = [
    (EME,  "Backend",   "FastAPI · Python 3.12 · asyncpg"),
    (SKY,  "Frontend",  "SvelteKit · TypeScript"),
    (AMB,  "AI Models", "GPT-4o + Claude Sonnet 4.6"),
    (VIO,  "Data",      "PostgreSQL · Redis · 12 tables"),
    (TEA,  "Infra",     "AWS RDS · ElastiCache · Secrets Manager"),
]
for i, (color, label, val) in enumerate(summary):
    x = 1.2 + i * 1.6
    box(sl, x, 5.5, 1.52, 0.82, S1, lc=BORD, lw=Pt(0.5))
    box(sl, x, 5.5, 1.52, 0.04, color)
    t(sl, label, x + 0.06, 5.55, 1.4, 0.25, color=color, size=9, bold=True)
    t(sl, val,   x + 0.06, 5.78, 1.4, 0.48, color=T2, size=8, wrap=True)

# Stack list in accent column
t(sl, "Architecture at a glance", 10.0, 0.7, 3.1, 0.3, color=GOLD, size=9, bold=True)
hline(sl, 10.0, 1.04, 3.05, color=BORD)
items = [
    "14 slides", "8 API routers", "7 agent tools",
    "12 DB tables", "4 sharing phases", "3 visibility tiers",
    "4 team roles", "1 permanent share URL",
]
for i, item in enumerate(items):
    box(sl, 10.0, 1.14 + i * 0.8, 0.04, 0.55, GOLD)
    t(sl, item, 10.12, 1.2 + i * 0.8, 3.0, 0.42, color=WT, size=11)


# ── Save ───────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Playbook-Research-Architecture.pptx")
prs.save(out)
print(f"✓  Saved  {out}  ({len(prs.slides)} slides)")
